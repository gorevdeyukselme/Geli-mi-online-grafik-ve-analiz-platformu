
# -*- coding: utf-8 -*-
"""
Gelişmiş Online Bilimsel Grafik Platformu
CSV, Excel ve NetCDF dosyalarından grafik, istatistik, harita, CTD profili ve rapor çıktısı üretir.
"""
from __future__ import annotations

import io
import re
import math
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st

st.set_page_config(page_title="Gelişmiş Online Grafik Platformu", page_icon="🌊", layout="wide")

# -------------------------------------------------------------------
# Genel yardımcılar
# -------------------------------------------------------------------
def clean_col(col) -> str:
    text = str(col).strip()
    text = re.sub(r"\s+", " ", text)
    return text


def normalize_numeric_text(series: pd.Series) -> pd.Series:
    s = series.astype(str).str.strip()
    s = s.replace(["", "nan", "NaN", "None", "NONE", "-", "—", "null", "NULL"], np.nan)
    s_tr = s.str.replace(".", "", regex=False).str.replace(",", ".", regex=False)
    numeric_tr = pd.to_numeric(s_tr, errors="coerce")
    numeric_en = pd.to_numeric(s, errors="coerce")
    return numeric_tr if numeric_tr.notna().sum() >= numeric_en.notna().sum() else numeric_en


def smart_convert(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out.columns = [clean_col(c) for c in out.columns]
    for col in out.columns:
        if out[col].dtype == "object":
            converted = normalize_numeric_text(out[col])
            if converted.notna().mean() >= 0.55:
                out[col] = converted
    return out


def read_csv_smart(uploaded_file) -> pd.DataFrame:
    raw = uploaded_file.getvalue()
    encodings = ["utf-8-sig", "utf-8", "cp1254", "latin1"]
    seps = [None, ";", ",", "\t"]
    last_err = None
    for enc in encodings:
        for sep in seps:
            try:
                return pd.read_csv(io.BytesIO(raw), encoding=enc, sep=sep, engine="python")
            except Exception as exc:
                last_err = exc
    raise ValueError(f"CSV/TXT dosyası okunamadı: {last_err}")


def read_excel_sheets(uploaded_file) -> Dict[str, pd.DataFrame]:
    raw = uploaded_file.getvalue()
    xl = pd.ExcelFile(io.BytesIO(raw))
    return {sheet: pd.read_excel(io.BytesIO(raw), sheet_name=sheet) for sheet in xl.sheet_names}


@st.cache_data(show_spinner=False)
def read_netcdf_bytes(raw: bytes, filename: str, max_rows: int) -> pd.DataFrame:
    import tempfile
    import xarray as xr
    suffix = Path(filename).suffix or ".nc"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name
    ds = xr.open_dataset(tmp_path)
    df = ds.to_dataframe().reset_index().dropna(axis=1, how="all")
    if len(df) > max_rows:
        df = df.head(max_rows)
    return df


def load_files(uploaded_files, netcdf_max_rows: int) -> Dict[str, pd.DataFrame]:
    datasets = {}
    for uf in uploaded_files:
        name = uf.name
        ext = Path(name).suffix.lower()
        try:
            if ext in [".csv", ".txt"]:
                datasets[name] = smart_convert(read_csv_smart(uf))
            elif ext in [".xlsx", ".xls"]:
                for sheet, df in read_excel_sheets(uf).items():
                    datasets[f"{name} | {sheet}"] = smart_convert(df)
            elif ext in [".nc", ".nc4", ".cdf", ".netcdf"]:
                datasets[name] = smart_convert(read_netcdf_bytes(uf.getvalue(), name, netcdf_max_rows))
            else:
                st.warning(f"Desteklenmeyen dosya türü: {name}")
        except Exception as exc:
            st.error(f"{name} okunamadı: {exc}")
    return datasets


@st.cache_data(show_spinner=False)
def load_demo_data() -> Dict[str, pd.DataFrame]:
    base = Path(__file__).parent / "sample_data"
    demo = {}
    for f in sorted(base.glob("*.csv")):
        demo[f"Örnek | {f.name}"] = smart_convert(pd.read_csv(f))
    return demo


def numeric_cols(df: pd.DataFrame) -> List[str]:
    return df.select_dtypes(include=[np.number]).columns.tolist()


def categorical_cols(df: pd.DataFrame, max_unique: int = 120) -> List[str]:
    cols = []
    for c in df.columns:
        if c not in numeric_cols(df) or df[c].nunique(dropna=True) <= max_unique:
            cols.append(c)
    return cols


def likely_datetime_cols(df: pd.DataFrame) -> List[str]:
    keys = ["date", "time", "tarih", "zaman", "datetime"]
    cols = []
    for c in df.columns:
        low = str(c).lower()
        if any(k in low for k in keys):
            cols.append(c)
            continue
        if df[c].dtype == "object":
            sample = df[c].dropna().astype(str).head(100)
            if len(sample) and pd.to_datetime(sample, errors="coerce", dayfirst=True).notna().mean() > 0.65:
                cols.append(c)
    return list(dict.fromkeys(cols))


def default_col(df: pd.DataFrame, keywords: List[str], only_numeric: bool = False) -> Optional[str]:
    cols = numeric_cols(df) if only_numeric else df.columns.tolist()
    for k in keywords:
        for c in cols:
            if k.lower() in str(c).lower():
                return c
    return cols[0] if cols else None


def order_categories_by_column(data: pd.DataFrame, category_col: str, order_col: Optional[str]) -> List:
    if category_col not in data.columns:
        return []
    if order_col and order_col != "Yok" and order_col in data.columns:
        order_series = data.groupby(category_col, dropna=False)[order_col].mean().sort_values()
        return order_series.index.tolist()
    try:
        return sorted(data[category_col].dropna().unique().tolist())
    except Exception:
        return data[category_col].dropna().astype(str).unique().tolist()


def add_contours_to_heatmap(fig, pivot: pd.DataFrame, contour_count: int = 10):
    try:
        z = pivot.astype(float).values
        if z.size < 4 or np.isfinite(z).sum() < 4:
            return fig
        zmin = float(np.nanmin(z))
        zmax = float(np.nanmax(z))
        if not np.isfinite(zmin) or not np.isfinite(zmax) or zmax <= zmin:
            return fig
        step = (zmax - zmin) / max(int(contour_count), 1)
        if step <= 0:
            return fig
        contour = go.Contour(
            z=z,
            x=list(pivot.columns),
            y=list(pivot.index),
            contours=dict(coloring="none", showlabels=True, start=zmin, end=zmax, size=step),
            line=dict(color="black", width=0.7),
            showscale=False,
            hoverinfo="skip",
        )
        fig.add_trace(contour)
    except Exception:
        pass
    return fig


def add_threshold_lines(fig, chart_type, x_threshold=None, y_threshold=None):
    if x_threshold is not None and chart_type != "3B saçılım grafiği":
        try:
            fig.add_vline(x=x_threshold, line_dash="dash", annotation_text=f"Eşik: {x_threshold}")
        except Exception:
            pass
    if y_threshold is not None and chart_type != "3B saçılım grafiği":
        try:
            fig.add_hline(y=y_threshold, line_dash="dash", annotation_text=f"Eşik: {y_threshold}")
        except Exception:
            pass
    return fig


def one_meter_average(df: pd.DataFrame, depth_col: str, group_cols: List[str], value_cols: List[str], start_at_one=True) -> pd.DataFrame:
    data = df.copy()
    data = data[pd.to_numeric(data[depth_col], errors="coerce").notna()]
    data = data[data[depth_col] >= 0]
    if start_at_one:
        data["Depth_bin_m"] = np.floor(data[depth_col]).astype(int)
        data.loc[data["Depth_bin_m"] < 1, "Depth_bin_m"] = 1
    else:
        data["Depth_bin_m"] = np.floor(data[depth_col]).astype(int)
    grouped = data.groupby(group_cols + ["Depth_bin_m"], dropna=False)[value_cols].mean().reset_index()
    return grouped


def qc_summary(df: pd.DataFrame) -> pd.DataFrame:
    rows = []
    nums = numeric_cols(df)
    for c in df.columns:
        item = {
            "Kolon": c,
            "Tip": str(df[c].dtype),
            "Eksik": int(df[c].isna().sum()),
            "Eksik_%": round(float(df[c].isna().mean()*100), 2),
            "Benzersiz": int(df[c].nunique(dropna=True)),
        }
        if c in nums:
            s = df[c].dropna()
            item.update({
                "Min": round(float(s.min()), 4) if len(s) else None,
                "Ort": round(float(s.mean()), 4) if len(s) else None,
                "Max": round(float(s.max()), 4) if len(s) else None,
                "Aykırı_IQR": int(((s < s.quantile(0.25)-1.5*(s.quantile(0.75)-s.quantile(0.25))) | (s > s.quantile(0.75)+1.5*(s.quantile(0.75)-s.quantile(0.25)))).sum()) if len(s)>3 else 0,
            })
        rows.append(item)
    return pd.DataFrame(rows)


def descriptive_stats(df: pd.DataFrame, group_col: Optional[str], selected_nums: List[str]) -> pd.DataFrame:
    if not selected_nums:
        return pd.DataFrame()
    if group_col and group_col != "Yok":
        return df.groupby(group_col)[selected_nums].agg(["count", "mean", "median", "std", "min", "max"]).round(4)
    return df[selected_nums].agg(["count", "mean", "median", "std", "min", "max"]).round(4)


def mann_kendall_test(y: pd.Series) -> Tuple[float, float, float]:
    vals = pd.to_numeric(y, errors="coerce").dropna().values
    n = len(vals)
    if n < 4:
        return np.nan, np.nan, np.nan
    s = 0
    for k in range(n-1):
        s += np.sign(vals[k+1:] - vals[k]).sum()
    var_s = n*(n-1)*(2*n+5)/18
    if s > 0:
        z = (s-1)/math.sqrt(var_s)
    elif s < 0:
        z = (s+1)/math.sqrt(var_s)
    else:
        z = 0.0
    # iki yönlü p için normal yaklaşım
    p = 2*(1 - 0.5*(1+math.erf(abs(z)/math.sqrt(2))))
    # Sen slope
    slopes = []
    for i in range(n-1):
        for j in range(i+1, n):
            slopes.append((vals[j]-vals[i])/(j-i))
    sen = float(np.median(slopes)) if slopes else np.nan
    return float(s), float(z), float(p), float(sen)


def pca_numpy(df: pd.DataFrame, cols: List[str]):
    X = df[cols].dropna().copy()
    if len(X) < 3 or len(cols) < 2:
        return None
    Xs = (X - X.mean()) / X.std(ddof=0).replace(0, np.nan)
    Xs = Xs.dropna()
    U, S, Vt = np.linalg.svd(Xs.values, full_matrices=False)
    scores = U[:, :2] * S[:2]
    loadings = Vt[:2, :].T
    exp = (S**2) / np.sum(S**2)
    return Xs.index, scores, loadings, exp[:2]


def fig_downloads(fig, base_name="grafik"):
    st.download_button("HTML indir", fig.to_html(include_plotlyjs="cdn").encode("utf-8"), f"{base_name}.html", "text/html")
    try:
        st.download_button("PNG indir", fig.to_image(format="png", scale=3), f"{base_name}.png", "image/png")
    except Exception:
        st.caption("PNG için kaleido paketi gerekir; requirements.txt içinde eklenmiştir. Online ortamda ilk çalıştırmada gecikebilir.")
    try:
        st.download_button("SVG indir", fig.to_image(format="svg"), f"{base_name}.svg", "image/svg+xml")
    except Exception:
        pass
    try:
        st.download_button("PDF indir", fig.to_image(format="pdf"), f"{base_name}.pdf", "application/pdf")
    except Exception:
        pass


def make_docx_report(title: str, summary_text: str, fig=None) -> bytes:
    from docx import Document
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph(summary_text)
    if fig is not None:
        try:
            img = fig.to_image(format="png", scale=2)
            bio = io.BytesIO(img)
            doc.add_picture(bio)
        except Exception:
            doc.add_paragraph("Grafik görseli eklenemedi; HTML/PNG indirme düğmesini kullanın.")
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def make_pptx_report(title: str, fig=None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if fig is not None:
        try:
            img = fig.to_image(format="png", scale=2)
            bio = io.BytesIO(img)
            slide.shapes.add_picture(bio, Inches(0.7), Inches(1.2), width=Inches(8.8))
        except Exception:
            tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            tx.text_frame.text = "Grafik görseli eklenemedi. HTML/PNG çıktısını ayrıca indirin."
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# -------------------------------------------------------------------
# Başlık ve yükleme
# -------------------------------------------------------------------
st.title("🌊 Gelişmiş Online Bilimsel Grafik Platformu")
st.caption("X/Y/Z eksen seçimi, CTD profili, harita, uydu grid, trend, PCA, korelasyon ve makale çıktıları.")

with st.sidebar:
    st.header("1) Veri")
    use_demo = st.checkbox("Örnek veri setlerini göster", value=True)
    uploaded_files = st.file_uploader(
        "CSV, Excel veya NetCDF yükle",
        type=["csv", "txt", "xlsx", "xls", "nc", "nc4", "cdf", "netcdf"],
        accept_multiple_files=True,
    )
    netcdf_max_rows = st.number_input("NetCDF azami satır", min_value=10_000, max_value=1_000_000, value=300_000, step=10_000)

all_datasets = {}
if use_demo:
    all_datasets.update(load_demo_data())
if uploaded_files:
    all_datasets.update(load_files(uploaded_files, netcdf_max_rows))

if not all_datasets:
    st.info("Başlamak için örnek veri setlerini açın veya dosya yükleyin.")
    st.stop()

with st.sidebar:
    dataset_name = st.selectbox("Veri seti", list(all_datasets.keys()))

raw = all_datasets[dataset_name].copy()

# -------------------------------------------------------------------
# Filtre ve temizlik
# -------------------------------------------------------------------
with st.sidebar:
    st.header("2) Temizleme / filtre")
    remove_empty_cols = st.checkbox("Tamamen boş kolonları kaldır", value=True)
    remove_duplicates = st.checkbox("Tekrarlı satırları kaldır", value=False)
    depth_guess = default_col(raw, ["Depth", "Derinlik"], only_numeric=True)
    depth_filter_col = st.selectbox("Negatif derinlik kontrol kolonu", ["Yok"] + numeric_cols(raw), index=(numeric_cols(raw).index(depth_guess)+1 if depth_guess in numeric_cols(raw) else 0))
    remove_negative_depth = st.checkbox("Negatif derinlikleri dışla", value=True)

    df = raw.copy()
    if remove_empty_cols:
        df = df.dropna(axis=1, how="all")
    if remove_duplicates:
        df = df.drop_duplicates()
    if remove_negative_depth and depth_filter_col != "Yok":
        df = df[df[depth_filter_col] >= 0]

    date_cols = likely_datetime_cols(df)
    date_col = st.selectbox("Tarih filtresi", ["Yok"] + date_cols)
    if date_col != "Yok":
        dt = pd.to_datetime(df[date_col], errors="coerce", dayfirst=True)
        if dt.notna().any():
            df["_date_filter_"] = dt
            mn, mx = dt.min().date(), dt.max().date()
            start_date, end_date = st.date_input("Tarih aralığı", value=(mn, mx))
            df = df[(df["_date_filter_"].dt.date >= start_date) & (df["_date_filter_"].dt.date <= end_date)]

    cat_f = st.selectbox("Kategori filtresi", ["Yok"] + categorical_cols(df))
    if cat_f != "Yok":
        vals = sorted([str(v) for v in df[cat_f].dropna().unique()])
        chosen = st.multiselect(f"{cat_f} seç", vals, default=vals[:min(30, len(vals))])
        if chosen:
            df = df[df[cat_f].astype(str).isin(chosen)]

st.subheader(f"Veri seti: {dataset_name}")
cols_metric = st.columns(5)
cols_metric[0].metric("Ham satır", f"{len(raw):,}")
cols_metric[1].metric("Aktif satır", f"{len(df):,}")
cols_metric[2].metric("Kolon", f"{df.shape[1]:,}")
cols_metric[3].metric("Sayısal kolon", f"{len(numeric_cols(df)):,}")
cols_metric[4].metric("Eksik hücre", f"{int(df.isna().sum().sum()):,}")

if df.empty:
    st.error("Filtrelerden sonra veri kalmadı.")
    st.stop()

nums = numeric_cols(df)
cols = df.columns.tolist()
if not nums:
    st.error("Sayısal kolon bulunamadı. Ondalık ayraçları kontrol edin.")
    st.stop()

tab_data, tab_general, tab_ctd, tab_map, tab_stats, tab_report = st.tabs([
    "📁 Veri & QC", "📊 Genel grafik", "🌊 CTD / profil", "🛰️ Harita / uydu", "📈 İstatistik", "📄 Rapor / çıktı"
])

# -------------------------------------------------------------------
# Veri & QC
# -------------------------------------------------------------------
with tab_data:
    st.markdown("### Veri önizleme")
    st.dataframe(df.head(1000), use_container_width=True)
    st.markdown("### Kalite kontrol özeti")
    qc = qc_summary(df)
    st.dataframe(qc, use_container_width=True)
    st.download_button("QC özetini CSV indir", qc.to_csv(index=False).encode("utf-8-sig"), "qc_ozeti.csv", "text/csv")

    st.markdown("### 1 m derinlik ortalaması")
    depth_col_1m = st.selectbox("Derinlik kolonu", nums, index=(nums.index(depth_guess) if depth_guess in nums else 0), key="depth_1m")
    val_cols_1m = st.multiselect("Ortalaması alınacak parametreler", nums, default=[c for c in nums if c != depth_col_1m][:min(6, max(1, len(nums)-1))])
    group_candidates = [c for c in categorical_cols(df, 300) if c != depth_col_1m]
    group_cols_1m = st.multiselect("Gruplama kolonları", group_candidates, default=[c for c in ["Station", "Cast_ID", "Date", "Month"] if c in group_candidates])
    start_one = st.checkbox("Derinliği 1 m'den başlat", value=True)
    if st.button("1 m ortalama tabloyu oluştur"):
        avg1m = one_meter_average(df, depth_col_1m, group_cols_1m, val_cols_1m, start_at_one=start_one)
        st.session_state["avg1m"] = avg1m
    if "avg1m" in st.session_state:
        st.dataframe(st.session_state["avg1m"].head(1000), use_container_width=True)
        st.download_button("1 m ortalama CSV indir", st.session_state["avg1m"].to_csv(index=False).encode("utf-8-sig"), "ctd_1m_ortalama.csv", "text/csv")

# -------------------------------------------------------------------
# Genel grafik
# -------------------------------------------------------------------
with tab_general:
    st.markdown("### Genel grafik çizici")
    c1, c2, c3, c4 = st.columns(4)
    chart_type = c1.selectbox("Grafik türü", [
        "Çizgi grafiği", "Saçılım grafiği", "3B saçılım grafiği", "Sütun grafiği", "Yatay sütun grafiği",
        "Kutu grafiği", "Violin plot", "Histogram", "Yoğunluk ısı haritası", "Yoğunluk kontur",
        "Pivot ısı haritası", "Radar grafiği", "Korelasyon matrisi"
    ])
    x_col = c2.selectbox("X", cols if chart_type not in ["Korelasyon matrisi"] else nums, key="g_x")
    y_col = c3.selectbox("Y", nums, index=min(1, len(nums)-1), key="g_y")
    z_col = c4.selectbox("Z / renk değeri", ["Yok"] + nums, key="g_z")
    color_col = st.selectbox("Renk/grup", ["Yok"] + cols, key="g_color")
    size_col = st.selectbox("Nokta boyutu", ["Yok"] + nums, key="g_size")
    title = st.text_input("Grafik başlığı", value=f"{chart_type} - {dataset_name}", key="g_title")
    trend = st.checkbox("Trend çizgisi", value=False, key="g_trend")
    x_thr_txt = st.text_input("X eşik çizgisi (boş bırakılabilir)", value="", key="g_xthr")
    y_thr_txt = st.text_input("Y eşik çizgisi (boş bırakılabilir)", value="", key="g_ythr")
    template = st.selectbox("Tema", ["plotly_white", "simple_white", "plotly", "ggplot2"], key="g_template")

    fig = None
    try:
        color_arg = None if color_col == "Yok" else color_col
        size_arg = None if size_col == "Yok" else size_col
        if chart_type == "Çizgi grafiği":
            data = df.dropna(subset=[x_col, y_col]).sort_values(x_col)
            fig = px.line(data, x=x_col, y=y_col, color=color_arg, markers=True, title=title, template=template)
        elif chart_type == "Saçılım grafiği":
            data = df.dropna(subset=[x_col, y_col])
            fig = px.scatter(data, x=x_col, y=y_col, color=color_arg, size=size_arg, trendline="ols" if trend else None, title=title, template=template)
        elif chart_type == "3B saçılım grafiği":
            z_axis = z_col if z_col != "Yok" else nums[min(2, len(nums)-1)]
            data = df.dropna(subset=[x_col, y_col, z_axis])
            fig = px.scatter_3d(data, x=x_col, y=y_col, z=z_axis, color=color_arg, size=size_arg, title=title, template=template)
        elif chart_type == "Sütun grafiği":
            fig = px.bar(df.dropna(subset=[x_col, y_col]), x=x_col, y=y_col, color=color_arg, barmode="group", title=title, template=template)
        elif chart_type == "Yatay sütun grafiği":
            fig = px.bar(df.dropna(subset=[x_col, y_col]), x=y_col, y=x_col, color=color_arg, orientation="h", title=title, template=template)
        elif chart_type == "Kutu grafiği":
            fig = px.box(df.dropna(subset=[x_col, y_col]), x=x_col, y=y_col, color=color_arg, points="all", title=title, template=template)
        elif chart_type == "Violin plot":
            fig = px.violin(df.dropna(subset=[x_col, y_col]), x=x_col, y=y_col, color=color_arg, box=True, points="all", title=title, template=template)
        elif chart_type == "Histogram":
            fig = px.histogram(df.dropna(subset=[x_col]), x=x_col, color=color_arg, marginal="box", title=title, template=template)
        elif chart_type == "Yoğunluk ısı haritası":
            fig = px.density_heatmap(df.dropna(subset=[x_col, y_col]), x=x_col, y=y_col, title=title, template=template)
        elif chart_type == "Yoğunluk kontur":
            fig = px.density_contour(df.dropna(subset=[x_col, y_col]), x=x_col, y=y_col, color=color_arg, title=title, template=template)
        elif chart_type == "Pivot ısı haritası":
            zval = z_col if z_col != "Yok" else y_col
            pivot = df.pivot_table(index=y_col, columns=x_col, values=zval, aggfunc="mean")
            fig = px.imshow(pivot, aspect="auto", title=title, template=template, labels=dict(x=x_col, y=y_col, color=zval))
        elif chart_type == "Radar grafiği":
            group = color_col if color_col != "Yok" else x_col
            radar_cols = [c for c in [y_col] + ([z_col] if z_col != "Yok" else []) if c in nums]
            if len(radar_cols) < 2:
                radar_cols = nums[:min(6, len(nums))]
            data = df.groupby(group)[radar_cols].mean().reset_index().head(8)
            fig = go.Figure()
            for _, row in data.iterrows():
                fig.add_trace(go.Scatterpolar(r=row[radar_cols].values, theta=radar_cols, fill="toself", name=str(row[group])))
            fig.update_layout(title=title, polar=dict(radialaxis=dict(visible=True)), template=template)
        elif chart_type == "Korelasyon matrisi":
            selected = st.multiselect("Korelasyona dahil edilecek değişkenler", nums, default=nums[:min(8,len(nums))], key="g_corr")
            corr = df[selected].corr()
            fig = px.imshow(corr, text_auto=".2f", zmin=-1, zmax=1, aspect="auto", title=title, template=template)
        def _parse_threshold(txt):
            try:
                txt = str(txt).strip()
                return None if txt == "" else float(txt.replace(",", "."))
            except Exception:
                return None
        fig = add_threshold_lines(fig, chart_type, _parse_threshold(x_thr_txt), _parse_threshold(y_thr_txt))
        fig.update_layout(height=720, title_x=0.02)
        st.plotly_chart(fig, use_container_width=True, key="general_plot")
        st.session_state["last_fig"] = fig
        st.session_state["last_title"] = title
        fig_downloads(fig, "genel_grafik")
    except Exception as exc:
        st.error(f"Grafik oluşturulamadı: {exc}")

# -------------------------------------------------------------------
# CTD / profil
# -------------------------------------------------------------------
with tab_ctd:
    st.markdown("### CTD ve dikey profil grafikleri")
    depth_default = default_col(df, ["Depth", "Derinlik"], only_numeric=True)
    temp_default = default_col(df, ["Temperature", "Temp", "Sıcak"], only_numeric=True)
    sal_default = default_col(df, ["Salinity", "Tuzluluk"], only_numeric=True)
    oxy_default = default_col(df, ["Oxygen", "Oksijen", "DO"], only_numeric=True)
    section_default = default_col(df, ["Transect", "Section", "Kesit", "Zone"], only_numeric=False)
    order_default = default_col(df, ["Distance_km", "Distance_mile", "Station_Order", "Order"], only_numeric=True)
    c1, c2, c3 = st.columns(3)
    ctd_chart = c1.selectbox("CTD grafik türü", ["Derinlik profili", "Çoklu profil", "T-S diyagramı", "Hovmöller", "Transekt / kesit", "Oksijen eşiği profili"])
    depth_col = c2.selectbox("Derinlik", nums, index=(nums.index(depth_default) if depth_default in nums else 0), key="ctd_depth")
    param_col = c3.selectbox("Parametre", nums, index=(nums.index(temp_default) if temp_default in nums else 0), key="ctd_param")
    group_col = st.selectbox("Renk/grup", ["Yok"] + cols, index=(cols.index("Station")+1 if "Station" in cols else 0), key="ctd_group")
    x_var = st.selectbox("Hovmöller/transekt X kolonu", cols, index=(cols.index("Month") if "Month" in cols else 0), key="ctd_xvar")

    extra1, extra2, extra3, extra4 = st.columns(4)
    section_col = extra1.selectbox("Kesit/Transect sütunu", ["Yok"] + categorical_cols(df), index=((["Yok"] + categorical_cols(df)).index(section_default) if section_default in categorical_cols(df) else 0), key="ctd_section_col")
    if section_col != "Yok":
        section_values = [str(v) for v in df[section_col].dropna().astype(str).unique().tolist()]
        selected_section = extra2.selectbox("Kesit seç", ["Tümü"] + section_values, key="ctd_section_val")
    else:
        selected_section = "Tümü"
        extra2.caption("İsteğe bağlı")
    order_col_ctd = extra3.selectbox("X sıralama kolonu", ["Yok"] + nums, index=((["Yok"] + nums).index(order_default) if order_default in nums else 0), key="ctd_order_col")
    show_contours = extra4.checkbox("Kontur/kesit çizgisi", value=(ctd_chart in ["Hovmöller", "Transekt / kesit"]), key="ctd_contours")
    contour_count = st.slider("Kontur çizgisi sayısı", min_value=4, max_value=20, value=10, key="ctd_contour_count") if show_contours and ctd_chart in ["Hovmöller", "Transekt / kesit"] else 10

    title_ctd = st.text_input("Başlık", value=f"{ctd_chart} - {param_col}", key="ctd_title")
    reverse_depth = st.checkbox("Derinlik aşağı doğru artsın", value=True, key="ctd_rev")
    fig_ctd = None
    try:
        color_arg = None if group_col == "Yok" else group_col
        data = df.copy()
        if section_col != "Yok" and selected_section != "Tümü":
            data = data[data[section_col].astype(str) == str(selected_section)]

        if ctd_chart in ["Derinlik profili", "Çoklu profil"]:
            data = data.dropna(subset=[depth_col, param_col]).sort_values(depth_col)
            fig_ctd = px.line(data, x=param_col, y=depth_col, color=color_arg, markers=True, title=title_ctd, template="plotly_white")
            if reverse_depth:
                fig_ctd.update_yaxes(autorange="reversed")
        elif ctd_chart == "Oksijen eşiği profili":
            oxy_col = st.selectbox("Oksijen kolonu", nums, index=(nums.index(oxy_default) if oxy_default in nums else 0), key="oxy_col")
            data = data.dropna(subset=[depth_col, oxy_col]).sort_values(depth_col)
            fig_ctd = px.line(data, x=oxy_col, y=depth_col, color=color_arg, markers=True, title=title_ctd, template="plotly_white")
            fig_ctd.add_vline(x=5, line_dash="dash", annotation_text="5 mg/L")
            fig_ctd.add_vline(x=2, line_dash="dot", annotation_text="2 mg/L")
            if reverse_depth:
                fig_ctd.update_yaxes(autorange="reversed")
        elif ctd_chart == "T-S diyagramı":
            sal_col = st.selectbox("Tuzluluk kolonu", nums, index=(nums.index(sal_default) if sal_default in nums else 0), key="ts_sal")
            temp_col = st.selectbox("Sıcaklık kolonu", nums, index=(nums.index(temp_default) if temp_default in nums else 0), key="ts_temp")
            color_ts = st.selectbox("Renk değeri", nums, index=(nums.index(oxy_default) if oxy_default in nums else 0), key="ts_color")
            data = data.dropna(subset=[sal_col, temp_col, color_ts])
            fig_ctd = px.scatter(data, x=sal_col, y=temp_col, color=color_ts, hover_data=[c for c in ["Station", "Depth_m", "Date", section_col if section_col != "Yok" else None] if c and c in cols], title=title_ctd, template="plotly_white")
        elif ctd_chart in ["Hovmöller", "Transekt / kesit"]:
            data = data.dropna(subset=[x_var, depth_col, param_col])
            pivot = data.pivot_table(index=depth_col, columns=x_var, values=param_col, aggfunc="mean")
            pivot = pivot.sort_index()
            if order_col_ctd != "Yok" and x_var in data.columns and order_col_ctd in data.columns:
                x_order = order_categories_by_column(data, x_var, order_col_ctd)
                pivot = pivot.reindex(columns=x_order)
            fig_ctd = px.imshow(pivot, aspect="auto", title=title_ctd, template="plotly_white", labels=dict(x=x_var, y=depth_col, color=param_col))
            if show_contours:
                fig_ctd = add_contours_to_heatmap(fig_ctd, pivot, contour_count=contour_count)
            if reverse_depth:
                fig_ctd.update_yaxes(autorange="reversed")

        fig_ctd.update_layout(height=720, title_x=0.02)
        st.plotly_chart(fig_ctd, use_container_width=True, key="ctd_plot")
        st.session_state["last_fig"] = fig_ctd
        st.session_state["last_title"] = title_ctd
        fig_downloads(fig_ctd, "ctd_grafik")
    except Exception as exc:
        st.error(f"CTD grafiği oluşturulamadı: {exc}")

# -------------------------------------------------------------------
# Harita / uydu
# -------------------------------------------------------------------
with tab_map:
    st.markdown("### Harita ve uydu/grid görselleştirme")
    lat_def = default_col(df, ["Latitude", "lat", "enlem"], only_numeric=True)
    lon_def = default_col(df, ["Longitude", "lon", "boylam"], only_numeric=True)
    label_default = default_col(df, ["Station", "İstasyon", "Name", "Label"], only_numeric=False)
    line_group_default = default_col(df, ["Transect", "Section", "Kesit", "Zone"], only_numeric=False)
    order_default = default_col(df, ["Station_Order", "Distance_km", "Distance_mile", "Order"], only_numeric=True)
    if lat_def is None or lon_def is None:
        st.warning("Harita için Latitude/Longitude kolonları gerekir.")
    else:
        c1, c2, c3, c4 = st.columns(4)
        map_type = c1.selectbox("Harita türü", ["İstasyon haritası", "Renkli nokta haritası", "Uydu grid ısı haritası", "Animasyonlu nokta haritası"])
        lat_col = c2.selectbox("Enlem", nums, index=nums.index(lat_def), key="lat")
        lon_col = c3.selectbox("Boylam", nums, index=nums.index(lon_def), key="lon")
        value_col = c4.selectbox("Renk/Z", nums, index=min(2, len(nums)-1), key="mapval")

        c5, c6, c7, c8 = st.columns(4)
        size_map = c5.selectbox("Boyut", ["Yok"] + nums, key="mapsize")
        label_col = c6.selectbox("İstasyon etiketi", ["Yok"] + cols, index=((["Yok"] + cols).index(label_default) if label_default in cols else 0), key="maplabel")
        show_labels = c7.checkbox("İstasyon adlarını göster", value=(label_default is not None), key="map_show_labels")
        show_lines = c8.checkbox("Kesit/transekt çizgisi", value=False, key="map_show_lines")

        c9, c10, c11 = st.columns(3)
        line_group_col = c9.selectbox("Kesit grubu", ["Yok"] + categorical_cols(df), index=((["Yok"] + categorical_cols(df)).index(line_group_default) if line_group_default in categorical_cols(df) else 0), key="maplinegroup")
        line_order_col = c10.selectbox("Çizgi sıralama kolonu", ["Yok"] + nums, index=((["Yok"] + nums).index(order_default) if order_default in nums else 0), key="maplineorder")
        anim_candidates = [c for c in cols if c != value_col]
        anim_col = c11.selectbox("Animasyon kolonu", ["Yok"] + anim_candidates, index=((["Yok"] + anim_candidates).index("Month") if "Month" in anim_candidates else 0), key="anim_col")

        hover_cols = [c for c in ["Station", label_col if label_col != "Yok" else None, "Transect", "Date", "Month", "Season", "Depth_m"] if c and c in cols]
        try:
            if map_type == "Uydu grid ısı haritası":
                data = df.dropna(subset=[lat_col, lon_col, value_col])
                pivot = data.pivot_table(index=lat_col, columns=lon_col, values=value_col, aggfunc="mean")
                fig_map = px.imshow(pivot, aspect="auto", origin="lower", title=f"Uydu/Grid haritası - {value_col}", labels=dict(x=lon_col, y=lat_col, color=value_col), template="plotly_white")
            elif map_type == "Animasyonlu nokta haritası" and anim_col != "Yok":
                data = df.dropna(subset=[lat_col, lon_col, value_col])
                fig_map = px.scatter_mapbox(data, lat=lat_col, lon=lon_col, color=value_col, size=None if size_map=="Yok" else size_map, animation_frame=anim_col, hover_data=hover_cols, text=label_col if show_labels and label_col != "Yok" else None, zoom=6, height=720, mapbox_style="open-street-map", title=f"Animasyonlu harita - {value_col}")
            else:
                data = df.dropna(subset=[lat_col, lon_col])
                fig_map = px.scatter_mapbox(data, lat=lat_col, lon=lon_col, color=None if map_type=="İstasyon haritası" else value_col, size=None if size_map=="Yok" else size_map, hover_data=hover_cols, text=label_col if show_labels and label_col != "Yok" else None, zoom=6, height=720, mapbox_style="open-street-map", title=f"Harita - {dataset_name}")

            if show_labels and map_type != "Uydu grid ısı haritası":
                fig_map.update_traces(textposition="top right")

            if show_lines and map_type != "Uydu grid ısı haritası":
                line_data = data.dropna(subset=[lat_col, lon_col]).copy()
                if line_group_col != "Yok":
                    grouped = line_data.groupby(line_group_col, dropna=False)
                else:
                    grouped = [("Tümü", line_data)]
                for grp_name, grp_df in grouped:
                    grp_df = grp_df.copy()
                    if line_order_col != "Yok" and line_order_col in grp_df.columns:
                        grp_df = grp_df.sort_values(line_order_col)
                    fig_map.add_trace(go.Scattermapbox(
                        lat=grp_df[lat_col],
                        lon=grp_df[lon_col],
                        mode="lines",
                        name=f"Kesit: {grp_name}",
                        line=dict(width=2),
                        hoverinfo="skip",
                        showlegend=True,
                    ))

            fig_map.update_layout(title_x=0.02, margin=dict(l=20,r=20,t=60,b=20))
            st.plotly_chart(fig_map, use_container_width=True, key="map_plot")
            st.session_state["last_fig"] = fig_map
            st.session_state["last_title"] = f"Harita - {dataset_name}"
            fig_downloads(fig_map, "harita_grafik")
        except Exception as exc:
            st.error(f"Harita oluşturulamadı: {exc}")

# -------------------------------------------------------------------
# İstatistik
# -------------------------------------------------------------------
with tab_stats:
    st.markdown("### İstatistik ve çok değişkenli analiz")
    stat_type = st.selectbox("Analiz", ["Tanımlayıcı istatistik", "Korelasyon", "PCA biplot", "Mann-Kendall trend", "Regresyon saçılım"])
    selected_nums = st.multiselect("Sayısal değişkenler", nums, default=nums[:min(6, len(nums))], key="stats_nums")
    group_col_stat = st.selectbox("Grup kolonu", ["Yok"] + categorical_cols(df), key="stat_group")
    try:
        if stat_type == "Tanımlayıcı istatistik":
            stat_df = descriptive_stats(df, group_col_stat, selected_nums)
            st.dataframe(stat_df, use_container_width=True)
            st.download_button("İstatistik CSV indir", stat_df.to_csv().encode("utf-8-sig"), "tanimlayici_istatistik.csv", "text/csv")
        elif stat_type == "Korelasyon":
            corr = df[selected_nums].corr(method=st.selectbox("Yöntem", ["pearson", "spearman", "kendall"]))
            fig_stat = px.imshow(corr, text_auto=".2f", zmin=-1, zmax=1, aspect="auto", title="Korelasyon matrisi", template="plotly_white")
            st.plotly_chart(fig_stat, use_container_width=True, key="stats_plot")
            st.session_state["last_fig"] = fig_stat
            st.session_state["last_title"] = "Korelasyon matrisi"
            fig_downloads(fig_stat, "korelasyon")
        elif stat_type == "PCA biplot":
            res = pca_numpy(df, selected_nums)
            if res is None:
                st.warning("PCA için en az 2 sayısal değişken ve yeterli veri gerekir.")
            else:
                idx, scores, loadings, exp = res
                score_df = pd.DataFrame({"PC1": scores[:,0], "PC2": scores[:,1]})
                if group_col_stat != "Yok":
                    score_df[group_col_stat] = df.loc[idx, group_col_stat].astype(str).values
                fig_stat = px.scatter(score_df, x="PC1", y="PC2", color=None if group_col_stat=="Yok" else group_col_stat, title=f"PCA biplot | PC1 {exp[0]*100:.1f}% - PC2 {exp[1]*100:.1f}%", template="plotly_white")
                scale = max(np.ptp(scores[:,0]), np.ptp(scores[:,1])) * 0.22
                for i, var in enumerate(selected_nums):
                    fig_stat.add_annotation(x=loadings[i,0]*scale, y=loadings[i,1]*scale, ax=0, ay=0, xref="x", yref="y", axref="x", ayref="y", text=var, showarrow=True, arrowhead=2)
                st.plotly_chart(fig_stat, use_container_width=True, key="stats_plot")
                st.session_state["last_fig"] = fig_stat
                st.session_state["last_title"] = "PCA biplot"
                fig_downloads(fig_stat, "pca_biplot")
        elif stat_type == "Mann-Kendall trend":
            time_col = st.selectbox("Zaman/X kolonu", cols, index=(cols.index("Year") if "Year" in cols else 0))
            ytrend = st.selectbox("Trend değişkeni", nums, index=min(0, len(nums)-1))
            data = df.dropna(subset=[time_col, ytrend]).sort_values(time_col)
            if data[time_col].duplicated().any():
                data = data.groupby(time_col, dropna=False)[ytrend].mean().reset_index()
            s, z, p, sen = mann_kendall_test(data[ytrend])
            fig_stat = px.scatter(data, x=time_col, y=ytrend, trendline="ols", title=f"Mann-Kendall trend | S={s:.1f}, Z={z:.2f}, p={p:.3f}, Sen slope={sen:.4f}", template="plotly_white")
            st.plotly_chart(fig_stat, use_container_width=True, key="stats_plot")
            st.info(f"Mann-Kendall: S={s:.1f}, Z={z:.2f}, p={p:.3f}; Sen eğimi={sen:.4f} birim/adım")
            st.session_state["last_fig"] = fig_stat
            st.session_state["last_title"] = "Mann-Kendall trend"
            fig_downloads(fig_stat, "trend_mann_kendall")
        elif stat_type == "Regresyon saçılım":
            xreg = st.selectbox("X", nums, index=0, key="xreg")
            yreg = st.selectbox("Y", nums, index=min(1, len(nums)-1), key="yreg")
            fig_stat = px.scatter(df.dropna(subset=[xreg, yreg]), x=xreg, y=yreg, color=None if group_col_stat=="Yok" else group_col_stat, trendline="ols", title=f"Regresyon: {yreg} ~ {xreg}", template="plotly_white")
            st.plotly_chart(fig_stat, use_container_width=True, key="stats_plot")
            st.session_state["last_fig"] = fig_stat
            st.session_state["last_title"] = f"Regresyon: {yreg} ~ {xreg}"
            fig_downloads(fig_stat, "regresyon")
    except Exception as exc:
        st.error(f"Analiz oluşturulamadı: {exc}")

# -------------------------------------------------------------------
# Rapor ve çıktı
# -------------------------------------------------------------------
with tab_report:
    st.markdown("### Makale / sunum çıktısı")
    last_fig = st.session_state.get("last_fig")
    last_title = st.session_state.get("last_title", "Grafik raporu")
    report_title = st.text_input("Rapor başlığı", value=last_title)
    summary_text = st.text_area("Rapor açıklaması", value=f"Veri seti: {dataset_name}\nAktif satır sayısı: {len(df):,}\nKolon sayısı: {df.shape[1]:,}")
    if last_fig is None:
        st.info("Önce herhangi bir sekmede grafik oluşturun. Son oluşturulan grafik burada rapora eklenir.")
    else:
        st.plotly_chart(last_fig, use_container_width=True, key="report_plot")
        fig_downloads(last_fig, "son_grafik")
        try:
            st.download_button("Word raporu indir", make_docx_report(report_title, summary_text, last_fig), "grafik_raporu.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        except Exception as exc:
            st.warning(f"Word raporu üretilemedi: {exc}")
        try:
            st.download_button("PowerPoint slaytı indir", make_pptx_report(report_title, last_fig), "grafik_slayti.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except Exception as exc:
            st.warning(f"PowerPoint üretilemedi: {exc}")
    st.download_button("Aktif veriyi CSV indir", df.to_csv(index=False).encode("utf-8-sig"), "aktif_veri.csv", "text/csv")
